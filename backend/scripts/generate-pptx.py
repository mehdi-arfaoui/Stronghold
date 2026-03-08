#!/usr/bin/env python3
"""
Fill Stronghold PPTX template using audit JSON data.
Usage:
  python3 generate-pptx.py --data data.json --template Template_Stronghold.pptx --output output.pptx
"""

import argparse
import json
from copy import deepcopy
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Pt


def load_data(path: str) -> dict[str, Any]:
    with open(path, "r", encoding="utf-8-sig") as handle:
        payload = json.load(handle)
    return payload if isinstance(payload, dict) else {}


def find_slide_by_title(prs: Presentation, title_contains: str):
    needle = (title_contains or "").strip().lower()
    if not needle:
        return None
    for slide in prs.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if isinstance(text, str) and needle in text.lower():
                return slide
    return None


def get_slide(prs: Presentation, index: int, title_hint: str | None = None):
    if 0 <= index < len(prs.slides):
        return prs.slides[index]
    if title_hint:
        return find_slide_by_title(prs, title_hint)
    return None


def find_shape_by_name(slide, name: str):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def set_text(shape, text: str, font_size: int | None = None, bold: bool | None = None, color: RGBColor | None = None):
    if not hasattr(shape, "text_frame"):
        return
    text_frame = shape.text_frame
    if not text_frame.paragraphs:
        paragraph = text_frame.add_paragraph()
    else:
        paragraph = text_frame.paragraphs[0]

    while len(paragraph.runs) > 1:
        paragraph._p.remove(paragraph.runs[-1]._r)

    if paragraph.runs:
        run = paragraph.runs[0]
    else:
        run = paragraph.add_run()

    run.text = str(text)
    if font_size is not None:
        run.font.size = Pt(font_size)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color

    while len(text_frame.paragraphs) > 1:
        text_frame._txBody.remove(text_frame.paragraphs[-1]._p)


def add_bullet_text(shape, lines: list[str], font_size: int = 11):
    if not hasattr(shape, "text_frame"):
        return
    text_frame = shape.text_frame
    while len(text_frame.paragraphs) > 1:
        text_frame._txBody.remove(text_frame.paragraphs[-1]._p)

    safe_lines = lines if lines else ["Aucune donnee disponible"]
    for index, line in enumerate(safe_lines):
        if index == 0:
            paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()
        paragraph.text = str(line)
        paragraph.level = 0
        paragraph.font.size = Pt(font_size)


def fill_table(table, headers: list[str], rows: list[list[Any]], header_color: RGBColor | None = None):
    if not headers:
        return

    for col_idx, header in enumerate(headers):
        if col_idx >= len(table.columns):
            break
        cell = table.cell(0, col_idx)
        cell.text = str(header)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(10)
            paragraph.font.bold = True
            if header_color is not None:
                paragraph.font.color.rgb = header_color

    safe_rows = rows if rows else [["Aucune donnee disponible"] + [""] * max(0, len(headers) - 1)]
    for row_idx, row_data in enumerate(safe_rows):
        actual_row = row_idx + 1
        if actual_row >= len(list(table.rows)):
            last_row = table._tbl.findall(qn("a:tr"))[-1]
            table._tbl.append(deepcopy(last_row))

        for col_idx, value in enumerate(row_data):
            if col_idx >= len(table.columns):
                break
            cell = table.cell(actual_row, col_idx)
            cell.text = "" if value is None else str(value)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)


def remove_slides(prs: Presentation, indexes: list[int]):
    for index in sorted(set(indexes), reverse=True):
        if 0 <= index < len(prs.slides):
            rel_id = prs.slides._sldIdLst[index].rId
            prs.part.drop_rel(rel_id)
            del prs.slides._sldIdLst[index]


def safe_percentage(value: Any) -> str:
    try:
        return f"{float(value):.0f}%"
    except Exception:
        return "0%"


def safe_currency(value: Any) -> str:
    if isinstance(value, (int, float)):
        return f"{value:,.0f} EUR"
    return str(value) if value is not None else "0 EUR"


def generate(data: dict[str, Any], template_path: str, output_path: str):
    prs = Presentation(template_path)

    # Slide 1 - Cover
    slide1 = get_slide(prs, 0, "Audit")
    if slide1:
        for shape in slide1.placeholders:
            if shape.placeholder_format.idx == 0:
                set_text(shape, "Audit de resilience IT", font_size=36, bold=True)
            elif shape.placeholder_format.idx == 1:
                client_name = str(data.get("clientName", "Client"))
                scan_date = str(data.get("scanDate", ""))
                set_text(shape, f"{client_name} - {scan_date}", font_size=16)

    # Slide 2 - Agenda
    slide2 = get_slide(prs, 1, "Agenda")
    if slide2:
        agenda_items = [
            "Contexte et objectifs",
            "Scenarios de risque",
            "Analyse de criticite",
            "Recommandations DR",
            "Plan d investissement DR",
            "Analyse BIA",
            "Feuille de route",
            "Actions en suspens",
        ]
        for shape in slide2.placeholders:
            if shape.placeholder_format.idx == 1:
                add_bullet_text(shape, agenda_items, font_size=14)

    # Slide 3 - divider
    slide3 = get_slide(prs, 2, "Context")
    if slide3:
        for shape in slide3.placeholders:
            if shape.placeholder_format.idx == 0:
                set_text(shape, "Contexte et objectifs")

    # Slide 4 - executive summary
    slide4 = get_slide(prs, 3, "Introduction")
    if slide4:
        for shape in slide4.placeholders:
            if shape.placeholder_format.idx == 0:
                set_text(shape, "Resume executif")
            elif shape.placeholder_format.idx == 2:
                score = data.get("resilienceScore", "N/A")
                total_nodes = int(data.get("totalNodes", 0) or 0)
                spof_count = int(data.get("spofCount", 0) or 0)
                reco_count = int(data.get("recommendationCount", 0) or 0)
                resilient_count = int(data.get("resilientByDesign", 0) or 0)
                providers = data.get("providers", [])
                dr_budget = data.get("totalDrCostAnnual", 0)
                provider_text = ", ".join([str(item) for item in providers]) if isinstance(providers, list) else ""

                lines = [
                    f"Score de resilience : {score}/100",
                    f"Services scannes : {total_nodes} ({provider_text})" if provider_text else f"Services scannes : {total_nodes}",
                    f"Services resilients par design : {resilient_count}",
                    f"SPOF identifies : {spof_count}",
                    f"Recommandations DR : {reco_count}",
                    f"Budget DR estime : {safe_currency(dr_budget)} / an",
                ]
                if bool(data.get("financialProfileConfigured", False)):
                    lines.append(f"ROI global : {data.get('globalRoi', 'N/A')}%")
                    lines.append(
                        f"Economie annuelle estimee : {safe_currency(data.get('totalAnnualSavings', 0))}"
                    )
                else:
                    lines.append("ROI : profil financier non configure")
                add_bullet_text(shape, lines, font_size=13)

    # Slide 5 divider
    slide5 = get_slide(prs, 4, "AWS DR Cases")
    if slide5:
        for shape in slide5.placeholders:
            if shape.placeholder_format.idx == 0:
                set_text(shape, "Scenarios de risque identifies")

    # Slide 6 - 3 columns for SPOF
    slide6 = get_slide(prs, 5, "Tableau")
    spofs = data.get("topSpofs", [])
    if not isinstance(spofs, list):
        spofs = []

    if slide6:
        title6 = find_shape_by_name(slide6, "Titre 1")
        if title6:
            set_text(title6, "SPOF identifies et impact")

        for index, shape_name in enumerate(["Rectangle 5", "Rectangle 6", "Rectangle 7"]):
            shape = find_shape_by_name(slide6, shape_name)
            if not shape:
                continue
            titles = ["SPOF detectes", "Impact (blast radius)", "Recommandation DR"]
            set_text(shape, titles[index])

        summary_shapes = ["Rectangle 14", "Rectangle 15", "Rectangle 16"]
        for index, shape_name in enumerate(summary_shapes):
            shape = find_shape_by_name(slide6, shape_name)
            if not shape:
                continue
            if index == 0:
                set_text(shape, f"{len(spofs)} SPOF critiques identifies")
            elif index == 1:
                avg_blast = 0
                if spofs:
                    avg_blast = sum(int(item.get("blastRadius", 0) or 0) for item in spofs) / len(spofs)
                set_text(shape, f"Blast radius moyen : {avg_blast:.0f} services")
            else:
                set_text(shape, f"{int(data.get('recommendationCount', 0) or 0)} actions recommandees")

        content_shapes = ["Rectangle 17", "Rectangle 18", "Rectangle 19"]
        recommendations = data.get("topRecommendations", [])
        if not isinstance(recommendations, list):
            recommendations = []
        for index, shape_name in enumerate(content_shapes):
            shape = find_shape_by_name(slide6, shape_name)
            if not shape:
                continue
            if index == 0:
                lines = [f"- {item.get('name', 'N/A')} ({item.get('type', '')})" for item in spofs[:8]]
                add_bullet_text(shape, lines or ["Aucun SPOF detecte"], font_size=10)
            elif index == 1:
                lines = [f"- {item.get('name', 'N/A')} : {item.get('blastRadius', 0)} services" for item in spofs[:8]]
                add_bullet_text(shape, lines or ["Aucun impact critique"], font_size=10)
            else:
                lines = [f"- {item.get('title', 'N/A')}" for item in recommendations[:8]]
                add_bullet_text(shape, lines or ["Aucune recommandation"], font_size=10)

    # Slide 7 divider
    slide7 = get_slide(prs, 6, "Criticality")
    if slide7:
        for shape in slide7.placeholders:
            if shape.placeholder_format.idx == 0:
                set_text(shape, "Analyse de criticite")

    # Slide 8 circles explanatory text
    slide8 = get_slide(prs, 7, "Criticality")
    if slide8:
        rect23 = find_shape_by_name(slide8, "Rectangle 23")
        if rect23:
            tiers = data.get("servicesByTier", {})
            if not isinstance(tiers, dict):
                tiers = {}
            text = "\n".join(
                [
                    f"Infrastructure analysee : {int(data.get('totalNodes', 0) or 0)} services.",
                    f"Tier 1 (critique) : {int(tiers.get('tier1', 0) or 0)} services",
                    f"Tier 2 (important) : {int(tiers.get('tier2', 0) or 0)} services",
                    f"Tier 3 (standard) : {int(tiers.get('tier3', 0) or 0)} services",
                    f"Tier 4 (non critique) : {int(tiers.get('tier4', 0) or 0)} services",
                ]
            )
            set_text(rect23, text, font_size=10)

    # Slide 9 divider
    slide9 = get_slide(prs, 8, "Recovery approach")
    if slide9:
        for shape in slide9.placeholders:
            if shape.placeholder_format.idx == 0:
                set_text(shape, "Recommandations DR")

    # Slide 10 recommendations title
    slide10 = get_slide(prs, 9, "Recovery")
    if slide10:
        title10 = find_shape_by_name(slide10, "Titre 1")
        if title10:
            set_text(title10, "Recommandations prioritaires (Quick Wins)")
        # Try to populate first available text shape.
        recommendations = data.get("topRecommendations", [])
        if not isinstance(recommendations, list):
            recommendations = []
        quick_wins = [item for item in recommendations if bool(item.get("quickWin", False))]
        if not quick_wins:
            quick_wins = recommendations[:5]
        lines = [
            f"- {item.get('title', 'N/A')} ({item.get('strategy', 'N/A')})"
            for item in quick_wins[:8]
        ] or ["Aucune recommandation prioritaire"]
        for shape in slide10.shapes:
            if hasattr(shape, "text_frame") and shape.name != "Titre 1":
                add_bullet_text(shape, lines, font_size=11)
                break

    # Slide 11 divider
    slide11 = get_slide(prs, 10, "Backup strategy")
    if slide11:
        for shape in slide11.placeholders:
            if shape.placeholder_format.idx == 0:
                set_text(shape, "Plan d investissement DR")

    # Slide 12 budget table
    slide12 = get_slide(prs, 11, "Backup strategy")
    if slide12:
        title12 = find_shape_by_name(slide12, "Titre 1")
        if title12:
            set_text(title12, "Repartition du budget DR par strategie")

        strategies = data.get("budgetByStrategy", [])
        if not isinstance(strategies, list):
            strategies = []
        rows = []
        for strategy in strategies:
            rows.append(
                [
                    strategy.get("strategy", ""),
                    str(strategy.get("count", 0)),
                    safe_currency(strategy.get("annualCost", 0)),
                    safe_percentage(strategy.get("percentage", 0)),
                ]
            )
        for shape in slide12.shapes:
            if shape.has_table:
                fill_table(
                    shape.table,
                    ["Strategie DR", "Nb recommandations", "Cout annuel", "% du budget"],
                    rows,
                )
                break

    # Slide 13 divider
    slide13 = get_slide(prs, 12, "Specific cases")
    if slide13:
        for shape in slide13.placeholders:
            if shape.placeholder_format.idx == 0:
                set_text(shape, "Analyse BIA")

    # Slide 14 BIA table
    slide14 = get_slide(prs, 13, "Specific cases")
    if slide14:
        title14 = find_shape_by_name(slide14, "Titre 1")
        if title14:
            set_text(title14, "Services critiques - RTO / RPO")

        bia_services = data.get("topBiaServices", [])
        if not isinstance(bia_services, list):
            bia_services = []
        rows = []
        for service in bia_services[:10]:
            rows.append(
                [
                    service.get("name", ""),
                    f"Tier {service.get('tier', '?')}",
                    service.get("rto", "N/A"),
                    service.get("rpo", "N/A"),
                ]
            )
        for shape in slide14.shapes:
            if shape.has_table:
                fill_table(shape.table, ["Service", "Tier", "RTO cible", "RPO cible"], rows)
                break

    # Slide 15 divider
    slide15 = get_slide(prs, 14, "Solutions")
    if slide15:
        for shape in slide15.placeholders:
            if shape.placeholder_format.idx == 0:
                set_text(shape, "Conformite reglementaire")

    # Slide 16 compliance table
    slide16 = get_slide(prs, 15, "Solutions")
    compliance = data.get("compliance", {})
    if not isinstance(compliance, dict):
        compliance = {}
    frameworks = compliance.get("frameworks", [])
    if not isinstance(frameworks, list):
        frameworks = []
    disclaimer = str(compliance.get("disclaimer", "")).strip()

    if slide16:
        title16 = find_shape_by_name(slide16, "Titre 1")
        if title16:
            set_text(title16, "Score de conformite")
        rows = []
        for framework in frameworks:
            rows.append(
                [
                    framework.get("name", ""),
                    f"{framework.get('score', 0)}%",
                    str(framework.get("compliant", 0)),
                    str(framework.get("partial", 0)),
                ]
            )
        if not rows:
            rows = [["ISO 22301", "Non evalue", "-", "-"], ["NIS 2", "Non evalue", "-", "-"]]
        for shape in slide16.shapes:
            if shape.has_table:
                fill_table(shape.table, ["Referentiel", "Score", "Conforme", "Partiel"], rows)
                break

        if disclaimer:
            for shape in slide16.shapes:
                if not hasattr(shape, "text_frame"):
                    continue
                if getattr(shape, "has_table", False):
                    continue
                if shape.name == "Titre 1":
                    continue
                set_text(shape, disclaimer, font_size=8, color=RGBColor(120, 120, 120))
                break

    # Slide 18 roadmap
    slide18 = get_slide(prs, 17, "Roadmap")
    if slide18:
        title18 = find_shape_by_name(slide18, "Titre 1")
        if title18:
            set_text(title18, "Feuille de route recommandee")

    # Slide 19 divider
    slide19 = get_slide(prs, 18, "Questions")
    if slide19:
        for shape in slide19.placeholders:
            if shape.placeholder_format.idx == 0:
                set_text(shape, "Actions en suspens")

    # Slide 20 pending actions
    slide20 = get_slide(prs, 19, "Questions")
    if slide20:
        title20 = find_shape_by_name(slide20, "Titre 1")
        if title20:
            set_text(title20, "Recommandations en attente de revue")
        pending = [
            item
            for item in (data.get("topRecommendations", []) if isinstance(data.get("topRecommendations", []), list) else [])
            if not bool(item.get("quickWin", False))
        ]
        lines = [
            f"- {item.get('title', 'N/A')} ({item.get('strategy', 'N/A')})"
            for item in pending[:8]
        ] or ["Aucune action en attente"]
        for shape in slide20.shapes:
            if hasattr(shape, "text_frame") and shape.name != "Titre 1":
                add_bullet_text(shape, lines, font_size=11)
                break

    # Optional compliance removal if no frameworks.
    slides_to_remove = []
    if len(frameworks) == 0:
        slides_to_remove.extend([14, 15])
    remove_slides(prs, slides_to_remove)

    prs.save(output_path)
    print(json.dumps({"success": True, "output": output_path, "slides": len(prs.slides)}))


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--data", required=True, help="Path to JSON data file")
    parser.add_argument("--template", required=True, help="Path to PPTX template")
    parser.add_argument("--output", required=True, help="Path to output PPTX file")
    args = parser.parse_args()

    data_path = Path(args.data)
    template_path = Path(args.template)
    output_path = Path(args.output)

    if not data_path.exists():
        raise FileNotFoundError(f"Data file not found: {data_path}")
    if not template_path.exists():
        raise FileNotFoundError(f"Template file not found: {template_path}")
    output_path.parent.mkdir(parents=True, exist_ok=True)

    payload = load_data(str(data_path))
    generate(payload, str(template_path), str(output_path))


if __name__ == "__main__":
    main()
